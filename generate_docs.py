"""
VelocityHire — Document Generator
Uses reportlab for PDFs, python-pptx for slides.
No system-level dependencies required.
"""

import os, re
from pathlib import Path
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.units import cm, mm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, HRFlowable,
    Table, TableStyle, PageBreak, Preformatted
)
from reportlab.pdfgen import canvas as rl_canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE       = Path("/mnt/efs/spaces/ba7ac73e-4169-4c04-979f-f044d06ab63e/e6e0512b-be61-46c7-9040-5f5821ef1dce")
MARKET_MD  = BASE / "velocityhire_marketing_plan.md"
PITCH_MD   = BASE / "velocityhire_pitch_deck.md"
MARKET_PDF = BASE / "velocityhire_marketing_plan.pdf"
PITCH_PDF  = BASE / "velocityhire_pitch_deck.pdf"
PITCH_PPTX = BASE / "velocityhire_pitch_deck.pptx"

# ── Colour palette ────────────────────────────────────────────────────────────
C_DARK    = colors.HexColor("#0D0D1A")
C_YELLOW  = colors.HexColor("#FFD700")
C_WHITE   = colors.white
C_LIGHT   = colors.HexColor("#F5F5FF")
C_MUTED   = colors.HexColor("#555577")
C_ACCENT  = colors.HexColor("#FFC107")
C_CODE_BG = colors.HexColor("#1E1E2E")
C_CODE_FG = colors.HexColor("#E0E0E0")

# ── PDF page decoration ───────────────────────────────────────────────────────
def make_canvas_decorator(title):
    def on_page(cv, doc):
        cv.saveState()
        w, h = A4
        # Top bar
        cv.setFillColor(C_DARK)
        cv.rect(0, h - 1.2*cm, w, 1.2*cm, fill=1, stroke=0)
        # Yellow accent line
        cv.setFillColor(C_YELLOW)
        cv.rect(0, h - 1.25*cm, w, 0.18*cm, fill=1, stroke=0)
        # Header text
        cv.setFillColor(C_YELLOW)
        cv.setFont("Helvetica-Bold", 9)
        cv.drawString(1.5*cm, h - 0.82*cm, "⚡ VELOCITYHIRE")
        cv.setFillColor(C_WHITE)
        cv.setFont("Helvetica", 8)
        cv.drawRightString(w - 1.5*cm, h - 0.82*cm, title)
        # Bottom bar
        cv.setFillColor(C_DARK)
        cv.rect(0, 0, w, 1.0*cm, fill=1, stroke=0)
        cv.setFillColor(C_YELLOW)
        cv.setFont("Helvetica-Bold", 7)
        cv.drawString(1.5*cm, 0.35*cm, "Complete.dev AI Agent Hackathon 2026")
        cv.setFillColor(colors.HexColor("#888899"))
        cv.setFont("Helvetica", 7)
        cv.drawRightString(w - 1.5*cm, 0.35*cm, f"Page {doc.page}")
        cv.restoreState()
    return on_page

# ── Build PDF styles ──────────────────────────────────────────────────────────
def get_styles():
    base = getSampleStyleSheet()
    s = {}
    def sty(name, **kw):
        s[name] = ParagraphStyle(name, **kw)

    sty("h1",    fontName="Helvetica-Bold", fontSize=22, textColor=C_DARK,
        spaceAfter=10, spaceBefore=18, leading=28,
        borderPad=(0,0,4,0), borderColor=C_YELLOW, borderWidth=0)
    sty("h2",    fontName="Helvetica-Bold", fontSize=15, textColor=C_DARK,
        spaceAfter=6, spaceBefore=16, leading=20,
        leftIndent=8, borderPad=(2,2,2,6), backColor=C_LIGHT)
    sty("h3",    fontName="Helvetica-Bold", fontSize=12, textColor=C_DARK,
        spaceAfter=4, spaceBefore=12, leading=16)
    sty("h4",    fontName="Helvetica-Bold", fontSize=10, textColor=C_MUTED,
        spaceAfter=4, spaceBefore=8, leading=14)
    sty("body",  fontName="Helvetica", fontSize=10, textColor=C_DARK,
        spaceAfter=4, leading=16)
    sty("bullet",fontName="Helvetica", fontSize=10, textColor=C_DARK,
        spaceAfter=3, leading=15, leftIndent=18, bulletIndent=6)
    sty("quote", fontName="Helvetica-Oblique", fontSize=11,
        textColor=C_WHITE, backColor=C_DARK,
        spaceAfter=10, spaceBefore=8, leading=17,
        leftIndent=10, rightIndent=10, borderPad=10,
        borderRadius=4)
    sty("code",  fontName="Courier", fontSize=8.5, textColor=C_CODE_FG,
        backColor=C_CODE_BG, spaceAfter=8, leading=13,
        leftIndent=10, rightIndent=10, borderPad=8)
    sty("tbl_h", fontName="Helvetica-Bold", fontSize=9.5, textColor=C_WHITE, leading=13)
    sty("tbl_c", fontName="Helvetica", fontSize=9, textColor=C_DARK, leading=13)
    return s

# ── Inline markdown cleaner ───────────────────────────────────────────────────
def clean_inline(text):
    # bold → <b>, code → <font> tag (reportlab XML)
    text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)
    text = re.sub(r'`(.*?)`', r'<font face="Courier" color="#c7254e" size="9">\1</font>', text)
    # strip markdown links but keep text
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    return text

# ── Parse a markdown block into reportlab flowables ──────────────────────────
def md_to_flowables(md_text, styles):
    elements = []
    lines = md_text.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]
        raw = line.rstrip()

        # Fenced code block
        if raw.startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].startswith('```'):
                code_lines.append(lines[i])
                i += 1
            code_text = '\n'.join(code_lines)
            elements.append(Preformatted(code_text, styles["code"]))
            i += 1
            continue

        # Heading 1
        if raw.startswith('# ') and not raw.startswith('## '):
            text = clean_inline(raw[2:].strip())
            elements.append(Spacer(1, 6))
            elements.append(Paragraph(text, styles["h1"]))
            elements.append(HRFlowable(width="100%", thickness=3, color=C_YELLOW, spaceAfter=6))
            i += 1; continue

        # Heading 2
        if raw.startswith('## '):
            text = clean_inline(raw[3:].strip())
            elements.append(Spacer(1, 4))
            elements.append(Paragraph(text, styles["h2"]))
            i += 1; continue

        # Heading 3
        if raw.startswith('### '):
            text = clean_inline(raw[4:].strip())
            elements.append(Paragraph(text, styles["h3"]))
            i += 1; continue

        # Heading 4
        if raw.startswith('#### '):
            text = clean_inline(raw[5:].strip())
            elements.append(Paragraph(text, styles["h4"]))
            i += 1; continue

        # Horizontal rule
        if raw.strip() in ('---', '***', '___'):
            elements.append(HRFlowable(width="100%", thickness=1, color=C_YELLOW, spaceBefore=8, spaceAfter=8))
            i += 1; continue

        # Blockquote
        if raw.startswith('> '):
            text = clean_inline(raw[2:].strip())
            elements.append(Paragraph(text, styles["quote"]))
            i += 1; continue

        # Table
        if raw.startswith('|'):
            table_rows = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                row_line = lines[i].strip()
                if re.match(r'^\|[-| :]+\|$', row_line):
                    i += 1; continue
                cells = [c.strip() for c in row_line.strip('|').split('|')]
                table_rows.append(cells)
                i += 1
            if table_rows:
                max_cols = max(len(r) for r in table_rows)
                for r in table_rows:
                    while len(r) < max_cols: r.append('')
                header = table_rows[0]
                body   = table_rows[1:]
                col_w  = [A4[0] / max_cols - 1.2*cm] * max_cols
                data   = [[Paragraph(clean_inline(c), styles["tbl_h"]) for c in header]]
                for row in body:
                    data.append([Paragraph(clean_inline(c), styles["tbl_c"]) for c in row])
                tbl = Table(data, colWidths=col_w, repeatRows=1)
                tbl.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), C_DARK),
                    ('TEXTCOLOR',  (0,0), (-1,0), C_YELLOW),
                    ('FONTNAME',   (0,0), (-1,0), 'Helvetica-Bold'),
                    ('FONTSIZE',   (0,0), (-1,0), 9.5),
                    ('ROWBACKGROUNDS', (0,1), (-1,-1), [C_WHITE, C_LIGHT]),
                    ('GRID',       (0,0), (-1,-1), 0.4, colors.HexColor("#D0D0E0")),
                    ('TOPPADDING', (0,0), (-1,-1), 5),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 5),
                    ('LEFTPADDING',  (0,0), (-1,-1), 7),
                    ('RIGHTPADDING', (0,0), (-1,-1), 7),
                    ('VALIGN',     (0,0), (-1,-1), 'MIDDLE'),
                ]))
                elements.append(tbl)
                elements.append(Spacer(1, 6))
            continue

        # Bullet list item
        if re.match(r'^[-*+]\s', raw) or re.match(r'^[→✅❌🏆⭐📋⚡📚🕐🌿⚙🎋]', raw):
            text = re.sub(r'^[-*+]\s+', '', raw)
            text = clean_inline(text)
            elements.append(Paragraph(f"• {text}", styles["bullet"]))
            i += 1; continue

        # Numbered list
        if re.match(r'^\d+\.\s', raw):
            text = re.sub(r'^\d+\.\s+', '', raw)
            text = clean_inline(text)
            elements.append(Paragraph(f"  {text}", styles["bullet"]))
            i += 1; continue

        # Empty line
        if not raw.strip():
            elements.append(Spacer(1, 4))
            i += 1; continue

        # Regular paragraph
        text = clean_inline(raw.strip())
        if text:
            elements.append(Paragraph(text, styles["body"]))
        i += 1

    return elements

def build_pdf(md_path: Path, pdf_path: Path, title: str):
    print(f"  Generating PDF: {pdf_path.name} ...")
    styles = get_styles()
    md_text = md_path.read_text()
    doc = SimpleDocTemplate(
        str(pdf_path),
        pagesize=A4,
        leftMargin=2.2*cm, rightMargin=2.2*cm,
        topMargin=2.2*cm, bottomMargin=1.8*cm,
        title=title,
        author="VelocityHire"
    )
    flowables = md_to_flowables(md_text, styles)
    doc.build(flowables, onFirstPage=make_canvas_decorator(title),
              onLaterPages=make_canvas_decorator(title))
    print(f"  ✅ PDF saved: {pdf_path.name}")


# ── PPTX Builder ──────────────────────────────────────────────────────────────
DARK_BG  = RGBColor(0x0D, 0x0D, 0x1A)
YELLOW_C = RGBColor(0xFF, 0xD7, 0x00)
WHITE_C  = RGBColor(0xFF, 0xFF, 0xFF)
GREY_C   = RGBColor(0xC0, 0xC5, 0xE0)
MUTED_C  = RGBColor(0x88, 0x88, 0xAA)

def set_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_rect(slide, prs, left, top, width, height, fill_rgb, line=False):
    from pptx.util import Inches
    shp = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if not line:
        shp.line.fill.background()
    return shp

def add_tb(slide, text, left, top, width, height,
           size=12, bold=False, color=WHITE_C, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def add_bullet_para(tf, text, size=11, color=GREY_C, level=1):
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    # strip markdown
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'`(.*?)`', r'\1', text)
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = "Calibri"

def build_pptx(pitch_md: Path, out_path: Path):
    print(f"  Generating PPTX: {out_path.name} ...")

    raw = pitch_md.read_text()
    # Split on slide dividers
    blocks = re.split(r'\n---\n', raw)
    slide_blocks = [b.strip() for b in blocks if re.match(r'##\s+SLIDE\s+\d+', b.strip(), re.IGNORECASE)]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for idx, block in enumerate(slide_blocks):
        lines = [l for l in block.split('\n') if l.strip()]

        # Slide title from "## SLIDE N — TITLE"
        hdr = lines[0] if lines else ""
        m = re.match(r'##\s+SLIDE\s+\d+\s*[—\-]+\s*(.*)', hdr, re.IGNORECASE)
        slide_title = m.group(1).strip() if m else f"Slide {idx+1}"

        content_lines = lines[1:]

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        for ph in slide.placeholders:
            ph._element.getparent().remove(ph._element)

        # Dark background
        set_bg(slide, 0x0D, 0x0D, 0x1A)

        # Top yellow bar
        add_rect(slide, prs, 0, 0, 13.33, 0.07, YELLOW_C)

        # Slide number
        add_tb(slide, f"{idx+1:02d}", 0.3, 0.12, 0.5, 0.35,
               size=9, bold=True, color=YELLOW_C)

        # Slide title
        add_tb(slide, slide_title.upper(), 0.55, 0.1, 12.5, 0.9,
               size=26, bold=True, color=YELLOW_C)

        # Yellow divider
        add_rect(slide, prs, 0.55, 1.0, 12.2, 0.04, YELLOW_C)

        # Content area textbox
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.12), Inches(12.2), Inches(5.9))
        tf = tb.text_frame
        tf.word_wrap = True
        first_para = True
        line_count = 0

        for line in content_lines:
            if line_count >= 24:
                break
            stripped = line.strip()
            if not stripped:
                continue

            # Clean inline markdown
            clean = re.sub(r'\*\*(.*?)\*\*', r'\1', stripped)
            clean = re.sub(r'`(.*?)`', r'\1', clean)
            clean = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', clean)
            clean = re.sub(r'^[-*+]\s+', '• ', clean)
            clean = re.sub(r'^\d+\.\s+', '  ', clean)

            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()

            # Detect line type
            is_heading3 = stripped.startswith('### ')
            is_heading4 = stripped.startswith('#### ')
            is_bullet   = clean.startswith('•') or any(stripped.startswith(s) for s in ['→', '✅', '❌', '🏆', '⭐', '📋', '⚡', '📚', '🕐', '🌿', '⚙', '🎋'])
            is_bold_kv  = re.match(r'\*\*.*?\*\*', stripped)
            is_quote    = stripped.startswith('> ')
            is_code     = stripped.startswith('    ') or stripped.startswith('```')

            run = p.add_run()

            if is_heading3 or is_heading4:
                clean = re.sub(r'^#+\s*', '', clean)
                run.text = clean
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = YELLOW_C
            elif is_quote:
                run.text = clean.lstrip('> ')
                run.font.size = Pt(12)
                run.font.italic = True
                run.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
            elif is_bullet:
                p.level = 1
                run.text = clean
                run.font.size = Pt(11)
                run.font.color.rgb = GREY_C
            elif is_bold_kv:
                run.text = clean
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = WHITE_C
            else:
                run.text = clean
                run.font.size = Pt(11)
                run.font.color.rgb = GREY_C

            run.font.name = "Calibri"
            line_count += 1

        # Footer
        add_tb(slide, "⚡ VelocityHire  ·  Complete.dev AI Agent Hackathon 2026  ·  https://q1inyxqs.run.complete.dev",
               0.5, 7.15, 12.3, 0.3,
               size=7.5, color=MUTED_C, align=PP_ALIGN.RIGHT)

    prs.save(str(out_path))
    print(f"  ✅ PPTX saved: {out_path.name}")


# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("\n⚡ VelocityHire Document Generator\n")

    print("📄 Generating PDFs...")
    build_pdf(MARKET_MD, MARKET_PDF, "VelocityHire — Marketing Plan")
    build_pdf(PITCH_MD,  PITCH_PDF,  "VelocityHire — Pitch Deck")

    print("\n📊 Generating PowerPoint slides...")
    build_pptx(PITCH_MD, PITCH_PPTX)

    print("\n✅ Done! Files generated:")
    for f in [MARKET_PDF, PITCH_PDF, PITCH_PPTX]:
        kb = f.stat().st_size // 1024
        print(f"   {f.name}  ({kb} KB)")
